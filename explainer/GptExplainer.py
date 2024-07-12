import asyncio
import json
import os
import logging
from pptx import Presentation
import openai
import pathlib

# Setup logging to output information to the console
logging.basicConfig(level=logging.INFO)

# Retrieve the OpenAI API key from environment variables
openai.api_key = os.getenv('OPENAI_API_KEY')
if not openai.api_key:
    raise ValueError("OpenAI API key not found. Please set the OPENAI_KEY environment variable.")

# Function to extract text from a slide
def extract_text_from_slide(slide):
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            texts.append(shape.text)
    return "\n".join(texts)

# Asynchronous function to get GPT explanation for the text
async def get_explanation(text):
    try:
        response = openai.Completion.create(
            model="gpt-3.5-turbo",
            prompt=f"Explain the following text:\n\n{text}\n\nExplanation:",
            max_tokens=150
        )
        explanation = response.choices[0].text.strip()
        return explanation
    except Exception as e:
        logging.error(f"Error getting explanation: {str(e)}")
        return -1  # Return -1 if an error occurs

# Asynchronous function to process the entire presentation
async def process_presentation(file_path):
    logging.info(f"Processing presentation: {file_path}")
    try:
        ppt = Presentation(file_path)
    except Exception as e:
        logging.error(f"Failed to open presentation: {str(e)}")
        return -1  # Return -1 if the presentation fails to open

    all_texts = [extract_text_from_slide(slide) for slide in ppt.slides]

    explanations = []
    for i, text in enumerate(all_texts):
        if text:
            logging.info(f"Submitting slide {i + 1}/{len(all_texts)} for explanation...")
            explanation = await get_explanation(text)
            if explanation == -1:
                return -1  # Return -1 if any slide processing fails
            explanations.append(explanation)
            if (i + 1) % 3 == 0:
                logging.info("Processed 3 slides. Waiting for 1 minute to comply with rate limits...")
                await asyncio.sleep(60)
            else:
                await asyncio.sleep(20)
    return explanations

# Function to save the explanations to a JSON file
def save_explanations(file_path, explanations):
    output_file = f"{os.path.splitext(file_path)[0]}.json"
    logging.info(f"Saving explanations to {output_file}")
    with open(output_file, 'w') as f:
        json.dump(explanations, f, indent=4)

# Main function to coordinate the presentation processing and saving
def MainFunc(path):
    file_path = pathlib.Path(path)
    explanations = asyncio.run(process_presentation(file_path))
    if explanations == -1:
        logging.error("No explanations were generated due to errors.")
        return -1
    save_explanations(file_path, explanations)
    logging.info(f"Explanations saved to {os.path.splitext(file_path)[0]}.json")
    return 0

if __name__ == "__main__":
    MainFunc(" ")
