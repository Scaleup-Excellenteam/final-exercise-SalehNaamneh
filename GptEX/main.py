import os
import json
import logging
import asyncio
from pptx import Presentation
import openai
from openai.error import APIError, APIConnectionError, RateLimitError, InvalidRequestError

# Set up logging
logging.basicConfig(level=logging.INFO)

# Load API key from environment variable
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")

if not OPENAI_API_KEY:
    logging.error("OPENAI_API_KEY is not set. Please set it in your environment variables.")
    exit(1)

# Set OpenAI API key
openai.api_key = OPENAI_API_KEY


def extract_text_from_slide(slide):
    """Extract text from a single slide."""
    # Joining text from all shapes that have text content in the slide
    return " ".join(shape.text for shape in slide.shapes if hasattr(shape, "text")).strip()


async def get_explanation(text, retries=5):
    """Get an explanation from OpenAI with exponential backoff and rate limiting."""
    delay = 1
    for attempt in range(retries):
        try:
            # Requesting explanation from OpenAI's GPT-3.5 model with specific parameters
            response = await openai.ChatCompletion.acreate(
                model="gpt-3.5-turbo",
                messages=[
                    {"role": "system", "content": "You are a helpful assistant."},
                    {"role": "user", "content": f"Explain the following presentation slide content succinctly:\n\n{text}"}
                ],
                max_tokens=1500,
                temperature=0.5,
            )
            # Logging the received explanation
            logging.info(f"Received explanation: {response.choices[0].message['content'].strip()[:60]}...")
            return response.choices[0].message['content'].strip()
        except RateLimitError:
            # Handling rate limit exceeded error by waiting for 1 minute before retrying
            logging.warning("Rate limit exceeded. Waiting for 1 minute before retrying...")
            await asyncio.sleep(60)
        except InvalidRequestError as e:
            # Handling insufficient quota error with a message and suggestion
            logging.error(f"Insufficient quota: {str(e)}")
            return "Insufficient quota to process this request. Please check your OpenAI plan and billing details."
        except (APIConnectionError, APIError) as e:
            # Handling generic OpenAI API errors and retrying with exponential backoff
            logging.error(f"OpenAI API error: {str(e)}. Retrying in {delay} seconds...")
            await asyncio.sleep(delay)
            delay = min(delay * 2, 60)
        except Exception as e:
            # Handling unexpected errors during API request or processing
            logging.error(f"Unexpected error: {str(e)}")
            return f"Error processing slide: {str(e)}"
    return "Failed to get explanation after several retries."


async def process_presentation(file_path):
    """Process the PowerPoint presentation asynchronously."""
    logging.info(f"Processing presentation: {file_path}")
    try:
        # Loading the PowerPoint presentation
        ppt = Presentation(file_path)
    except Exception as e:
        # Logging error if PowerPoint file cannot be opened
        logging.error(f"Failed to open presentation: {str(e)}")
        return []

    # Extracting text from each slide and preparing for explanation retrieval
    all_texts = [extract_text_from_slide(slide) for slide in ppt.slides]

    explanations = []
    for i, text in enumerate(all_texts):
        if text:
            # Submitting each slide's text for explanation generation
            logging.info(f"Submitting slide {i + 1}/{len(all_texts)} for explanation...")
            explanation = await get_explanation(text)
            explanations.append(explanation)
            if (i + 1) % 3 == 0:
                # Logging to comply with rate limits by waiting 1 minute after every 3 slides
                logging.info("Processed 3 slides. Waiting for 1 minute to comply with rate limits...")
                await asyncio.sleep(60)
            else:
                await asyncio.sleep(20)  # Ensuring minimum time between requests to avoid rate limits
    return explanations


def save_explanations(file_path, explanations):
    """Save explanations to a JSON file."""
    output_file = f"{os.path.splitext(file_path)[0]}.json"
    logging.info(f"Saving explanations to {output_file}")
    # Writing explanations to a JSON file with indentation for readability
    with open(output_file, 'w') as f:
        json.dump(explanations, f, indent=4)


def main():
    """Main function to process the PowerPoint file and save explanations."""
    file_path = "logging, debugging, getting into a large codebase.pptx"  # Replace with your PowerPoint file path this is my check PowerPoint file

    explanations = asyncio.run(process_presentation(file_path))
    if explanations:
        # Saving generated explanations to a JSON file
        save_explanations(file_path, explanations)
        logging.info(f"Explanations saved to {os.path.splitext(file_path)[0]}.json")
    else:
        logging.error("No explanations were generated due to errors.")


if __name__ == "__main__":
    main()

